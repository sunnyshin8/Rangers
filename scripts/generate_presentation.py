#!/usr/bin/env python3
"""Generate a 4-slide PowerPoint presentation for Leak Radar.
Requires: python-pptx

Usage:
  python scripts/generate_presentation.py --out leak-radar/presentation/leak_radar_presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import argparse

def make_presentation(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    assets = Path('leak-radar/assets')
    cover = assets / 'cover_cartoon.png'
    if not cover.exists():
        cover = assets / 'cover.png'

    # Slide 1 - Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if cover.exists():
        slide.shapes.add_picture(str(cover), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    else:
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = 'Leak Radar'
        p.font.size = Pt(64)
        p.font.bold = True

    # Slide 2 - Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = 'Problem'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = 'Secrets and sensitive data often leak into repos, CI logs, and public feeds — detection is slow and noisy.'
    p.level = 0
    for t in [
        'Developers lack prioritized, actionable findings.',
        'High false-positive rates waste security time.',
        'Remediation is manual and slow.'
    ]:
        p = body.add_paragraph()
        p.text = '• ' + t
        p.level = 1

    # Slide 3 - Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Solution: Leak Radar'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = 'Continuous scanning across repos, feeds, and storage with ML-based prioritization and automated policy-driven responses.'
    for t in [
        'Fast scanners + ML classifier to reduce false positives',
        'Policy engine for automated alerts, tickets, or remediation hooks',
        'Integrations: CI/CD, Slack, Jira, SIEMs and audit exports'
    ]:
        p = body.add_paragraph()
        p.text = '• ' + t
        p.level = 1

    # Slide 4 - Why Leak Radar is unique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Why Leak Radar is Unique'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = 'Cost-effective, buffer-free, and developer-friendly leak detection.'
    for t in [
        'Cost-effective: incremental scans, sampling, and prioritized findings reduce compute and analyst time',
        'Buffer-free system design: streaming architecture with backpressure and real-time processing (low latency)',
        'Privacy-first: configurable redaction and audit-ready evidence',
        'Integrated, animated monitoring: rich alerts (Slack/Teams/DMs) with interactive remediation actions'
    ]:
        p = body.add_paragraph()
        p.text = '• ' + t
        p.level = 1

    # Slide 5 - Call to action
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Next Steps & Demo'
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = 'Try a demo scan, tune detectors, and review prioritized findings.'
    p = body.add_paragraph()
    p.text = '• Demo: run sample scan against a test repo'
    p.level = 1
    p = body.add_paragraph()
    p.text = '• Contact: team@leakradar.example'
    p.level = 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print('Wrote', out_path)

if __name__ == '__main__':
    import sys
    parser = argparse.ArgumentParser()
    parser.add_argument('--out', default='leak-radar/presentation/leak_radar_presentation.pptx')
    args = parser.parse_args()
    make_presentation(Path(args.out))
